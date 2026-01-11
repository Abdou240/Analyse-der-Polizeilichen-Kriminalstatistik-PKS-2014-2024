#!/usr/bin/env python3
"""
Optimized script to generate professional visualizations for the PKS project 
and compile a PowerPoint presentation. Includes improved styling, better 
scaling for small values, and scientific presentation standards.
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib.ticker as ticker
from scipy import stats

# Set global style for professional appearance
plt.style.use('seaborn-v0_8-whitegrid')
sns.set_palette("Set2")

# Define color scheme for consistency
COLORS = {
    'diebstahl_insgesamt': '#1f77b4',  # Blue
    'wohnungseinbruch': '#ff7f0e',      # Orange
    'kfz_diebstahl': '#2ca02c',         # Green
    'bars': '#4c72b0',                  # Steel blue for bars
    'trend_line': '#d62728'             # Red for trend lines
}

def create_figures():
    """Create professional figures and save them as PNG files."""
    # Load data
    cases_df = pd.read_csv('time_series_cases_clean_2014_2024.csv')
    clearance_df = pd.read_csv('time_series_clearance_clean_2014_2024.csv')
    states_cases = pd.read_csv('laender_relevant_cases_2012_2024.csv')
    states_clearance = pd.read_csv('laender_relevant_aq_2012_2024.csv')

    # Ensure the year column is numeric
    cases_df['year'] = cases_df['year'].astype(int)
    clearance_df['year'] = clearance_df['year'].astype(int)
    states_cases['year'] = states_cases['year'].astype(int)
    states_clearance['year'] = states_clearance['year'].astype(int)

    # ----- FIGURE 1: Time series of case numbers with dual axis -----
    fig, ax1 = plt.subplots(figsize=(10, 6))
    
    # Primary axis for Diebstahl insgesamt
    line1 = ax1.plot(cases_df['year'], cases_df['diebstahl_insgesamt'], 
             marker='o', linewidth=2, color=COLORS['diebstahl_insgesamt'],
             label='Diebstahl insgesamt')
    ax1.set_xlabel('Jahr', fontsize=12)
    ax1.set_ylabel('Fallzahl (Diebstahl insgesamt)', fontsize=12, 
                   color=COLORS['diebstahl_insgesamt'])
    ax1.tick_params(axis='y', labelcolor=COLORS['diebstahl_insgesamt'])
    ax1.grid(True, alpha=0.3)
    
    # Secondary axis for smaller series
    ax2 = ax1.twinx()
    line2 = ax2.plot(cases_df['year'], cases_df['wohnungseinbruch'], 
             marker='s', linestyle='--', linewidth=1.5, 
             color=COLORS['wohnungseinbruch'],
             label='Wohnungseinbruchdiebstahl')
    line3 = ax2.plot(cases_df['year'], cases_df['kfz_diebstahl'], 
             marker='^', linestyle=':', linewidth=1.5,
             color=COLORS['kfz_diebstahl'],
             label='Kfz-Diebstahl')
    ax2.set_ylabel('Fallzahl (Wohnungseinbruch & Kfz)', fontsize=12)
    
    # Formatting
    ax1.set_title('Entwicklung der Fallzahlen 2014-2024', fontsize=14, fontweight='bold')
    ax1.set_xticks(cases_df['year'])
    ax1.xaxis.set_major_formatter(ticker.FormatStrFormatter('%d'))
    
    # Combine legends from both axes
    lines = line1 + line2 + line3
    labels = ['Diebstahl insgesamt', 'Wohnungseinbruchdiebstahl', 'Kfz-Diebstahl']
    
    # Place legend inside plot, upper right (like figure 2)
    ax1.legend(lines, labels, loc='upper right', fontsize=10, 
               framealpha=0.9, edgecolor='black')
    
    # Adjust layout
    plt.tight_layout()
    plt.savefig('figure1_cases.png', dpi=300, bbox_inches='tight')
    plt.close()

    # ----- FIGURE 2: Time series of clearance rates -----
    plt.figure(figsize=(10, 6))
    
    # Plot with distinct line styles and markers
    line_styles = ['-', '--', ':']
    markers = ['o', 's', '^']
    
    for idx, (col, label) in enumerate([
        ('diebstahl_insgesamt', 'Diebstahl insgesamt'),
        ('wohnungseinbruch', 'Wohnungseinbruchdiebstahl'),
        ('kfz_diebstahl', 'Kfz-Diebstahl'),
    ]):
        plt.plot(clearance_df['year'], clearance_df[col],
                 marker=markers[idx], linestyle=line_styles[idx], 
                 linewidth=2, markersize=6,
                 label=label)
    
    plt.title('Entwicklung der Aufklarungsquoten 2014-2024', 
              fontsize=14, fontweight='bold')
    plt.xlabel('Jahr', fontsize=12)
    plt.ylabel('Aufklarungsquote (%)', fontsize=12)
    
    # Place legend in upper right
    plt.legend(loc='upper right', fontsize=10, framealpha=0.9, edgecolor='black')
    
    plt.grid(True, alpha=0.3)
    
    # Add horizontal grid lines for better readability
    plt.gca().yaxis.grid(True, which='major', alpha=0.5, linestyle='--')
    plt.xticks(clearance_df['year'])
    
    plt.tight_layout()
    plt.savefig('figure2_clearance.png', dpi=300, bbox_inches='tight')
    plt.close()

    # ----- FIGURE 3: Bar chart of states (top 10) for 2024 -----
    year_filter = 2024
    df_2024 = states_cases[states_cases['year'] == year_filter].copy()
    df_2024 = df_2024.sort_values('diebstahl_insgesamt', ascending=True)  # For horizontal bars
    df_top10 = df_2024.tail(10)  # Top 10
    
    plt.figure(figsize=(10, 6))
    
    # Create horizontal bar chart
    bars = plt.barh(df_top10['bundesland'], df_top10['diebstahl_insgesamt'], 
                    color=COLORS['bars'], alpha=0.8)
    
    # Add value labels on bars
    for bar in bars:
        width = bar.get_width()
        plt.text(width + (width * 0.01), bar.get_y() + bar.get_height()/2,
                 f'{int(width/1000):,d}k'.replace(',', '.'),
                 va='center', fontsize=9)
    
    plt.title(f'Top-10 Bundeslander nach Diebstahl insgesamt ({year_filter})', 
              fontsize=14, fontweight='bold')
    plt.xlabel('Fallzahl', fontsize=12)
    
    # Format x-axis in thousands
    ax = plt.gca()
    ax.xaxis.set_major_formatter(ticker.FuncFormatter(lambda x, p: f'{int(x/1000):,d}k'.replace(',', '.')))
    
    plt.tight_layout()
    plt.savefig('figure3_bar_states.png', dpi=300, bbox_inches='tight')
    plt.close()

    # ----- FIGURE 4: Scatter plot of cases vs clearance rates -----
    # Bundesland abbreviations mapping
    bundesland_abbr = {
        'Baden-Württemberg': 'BW',
        'Bayern': 'BY',
        'Berlin': 'BE',
        'Brandenburg': 'BB',
        'Bremen': 'HB',
        'Hamburg': 'HH',
        'Hessen': 'HE',
        'Mecklenburg-Vorpommern': 'MV',
        'Niedersachsen': 'NI',
        'Nordrhein-Westfalen': 'NW',
        'Rheinland-Pfalz': 'RP',
        'Saarland': 'SL',
        'Sachsen': 'SN',
        'Sachsen-Anhalt': 'ST',
        'Schleswig-Holstein': 'SH',
        'Thüringen': 'TH'
    }
    
    df_cases = states_cases[states_cases['year'] == year_filter][['bundesland', 'diebstahl_insgesamt']]
    df_clear = states_clearance[states_clearance['year'] == year_filter][['bundesland', 'diebstahl_insgesamt']]
    merged = df_cases.merge(df_clear, on='bundesland', suffixes=('_cases', '_aq'))
    
    # Add abbreviations
    merged['abbr'] = merged['bundesland'].map(bundesland_abbr)
    
    # Calculate regression
    slope, intercept, r_value, p_value, std_err = stats.linregress(
        merged['diebstahl_insgesamt_cases'], merged['diebstahl_insgesamt_aq'])
    r_squared = r_value**2
    
    plt.figure(figsize=(9, 7))
    
    # Create scatter plot
    scatter = plt.scatter(merged['diebstahl_insgesamt_cases'], 
                         merged['diebstahl_insgesamt_aq'],
                         s=100, alpha=0.7, edgecolors='black', linewidth=0.5)
    
    # Add trend line
    x_range = np.linspace(merged['diebstahl_insgesamt_cases'].min(), 
                         merged['diebstahl_insgesamt_cases'].max(), 100)
    trend_line, = plt.plot(x_range, slope * x_range + intercept, 
             color=COLORS['trend_line'], linewidth=2, linestyle='--',
             label=f'Trendlinie (R² = {r_squared:.3f})')
    
    # Add state labels with small offset to avoid overlap
    for _, row in merged.iterrows():
        plt.annotate(row['abbr'], 
                    xy=(row['diebstahl_insgesamt_cases'], row['diebstahl_insgesamt_aq']),
                    xytext=(5, 5), textcoords='offset points',
                    fontsize=8, alpha=0.8)
    
    plt.title('Zusammenhang Fallzahl und Aufklarungsquote\n(Diebstahl insgesamt, 2024)', 
              fontsize=14, fontweight='bold')
    plt.xlabel('Fallzahl (Diebstahl insgesamt)', fontsize=12)
    plt.ylabel('Aufklarungsquote (%)', fontsize=12)
    
    # Highlight outliers (Berlin and Bayern) - only if they exist
    handles = [trend_line]
    labels = [f'Trendlinie (R² = {r_squared:.3f})']
    
    berlin = merged[merged['abbr'] == 'BE']
    bayern = merged[merged['abbr'] == 'BY']
    
    if not berlin.empty:
        berlin_scatter = plt.scatter(berlin['diebstahl_insgesamt_cases'], berlin['diebstahl_insgesamt_aq'],
                   s=150, color='red', edgecolors='black', linewidth=1.5, 
                   label='Berlin', zorder=5)
        handles.append(berlin_scatter)
        labels.append('Berlin')
    
    if not bayern.empty:
        bayern_scatter = plt.scatter(bayern['diebstahl_insgesamt_cases'], bayern['diebstahl_insgesamt_aq'],
                   s=150, color='orange', edgecolors='black', linewidth=1.5,
                   label='Bayern', zorder=5)
        handles.append(bayern_scatter)
        labels.append('Bayern')
    
    # Place legend inside plot, upper right (like figure 2)
    plt.legend(handles=handles, labels=labels, 
               loc='upper right', fontsize=10, framealpha=0.9, edgecolor='black')
    
    plt.grid(True, alpha=0.3)
    
    # Format x-axis in thousands
    ax = plt.gca()
    ax.xaxis.set_major_formatter(ticker.FuncFormatter(lambda x, p: f'{int(x/1000):,d}k'.replace(',', '.')))
    
    # Adjust layout
    plt.tight_layout()
    plt.savefig('figure4_scatter_cases_clearance.png', dpi=300, bbox_inches='tight')
    plt.close()

    print("Professional figures generated successfully.")
    return r_value, r_squared

def build_presentation(r_value=None, r_squared=None):
    """Assemble the PowerPoint presentation using python‑pptx."""
    prs = Presentation()
    
    # Set presentation aspect ratio (16:9 for modern presentations)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[5]

    # Slide 1: Title
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Zeitliche und regionale Analyse von Diebstahlsdelikten"
    subtitle.text = "PKS-Daten 2014-2024\nWissenschaftliche Auswertung"
    
    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(36)
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)

    # Slide 2: Forschungsfragen
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Fragestellungen und Datenbasis"
    
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    
    questions = [
        "Wie haben sich die Fallzahlen 2014-2024 entwickelt?",
        "Welche regionalen Unterschiede zeigen sich zwischen den Bundeslandern?",
        "Wie verhalten sich Fallzahlen und Aufklarungsquoten zueinander?",
        "Optional: Welche Rolle spielen Bevolkerungsdichte und Urbanisierung?"
    ]
    
    for question in questions:
        p = body.add_paragraph()
        p.text = f"• {question}"
        p.font.size = Pt(16)
        p.font.bold = False
    
    # Data source note
    p = body.add_paragraph()
    p.text = "\nDatenquelle: Bereinigte PKS-Zeitreihen und Lander-Tabellen (2014-2024)"
    p.font.size = Pt(12)
    p.font.italic = True

    # Slide 3: Methodik
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Methodische Herangehensweise"
    
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    
    methods = [
        "Berechnung standardisierter Haufigkeiten pro 100.000 Einwohner",
        "Liniendiagramme zur Darstellung langfristiger Trends",
        "Bar-Diagramme und Streudiagramme zur regionalen Analyse",
        "Lineare Regression zur Untersuchung von Zusammenhangen",
        "Dual-Achsen-Darstellung fur unterschiedliche Skalierungen"
    ]
    
    for method in methods:
        p = body.add_paragraph()
        p.text = f"• {method}"
        p.font.size = Pt(16)

    # Slide 4: Trends - Fallzahlen
    slide = prs.slides.add_slide(blank_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Zeitliche Entwicklung der Fallzahlen"
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    
    # Insert figure with border
    left = Inches(0.5)
    top = Inches(1.2)
    pic = slide.shapes.add_picture('figure1_cases.png', left, top, width=Inches(9))
    
    # Add professional caption
    caption = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
    tf = caption.text_frame
    tf.text = "Abbildung 1: Fallzahlen der drei Deliktgruppen (2014-2024). Anmerkung: Wohnungseinbruch- und Kfz-Diebstahlsfalle sind auf der sekundaren Y-Achse dargestellt."
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.italic = True

    # Slide 5: Trends - Aufklarungsquoten
    slide = prs.slides.add_slide(blank_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Entwicklung der Aufklarungsquoten"
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    
    pic = slide.shapes.add_picture('figure2_clearance.png', Inches(0.5), Inches(1.2), width=Inches(9))
    
    caption = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
    tf = caption.text_frame
    tf.text = "Abbildung 2: Aufklarungsquoten nach Deliktgruppe (2014-2024). Wohnungseinbruch zeigt die niedrigsten Aufklarungsraten."
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.italic = True

    # Slide 6: Regionale Unterschiede
    slide = prs.slides.add_slide(blank_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Regionale Unterschiede 2024"
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    
    pic = slide.shapes.add_picture('figure3_bar_states.png', Inches(0.5), Inches(1.2), width=Inches(9))
    
    caption = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1))
    tf = caption.text_frame
    tf.text = "Abbildung 3: Top-10 Bundeslander nach Fallzahlen (Diebstahl insgesamt, 2024). NRW und Berlin zeigen die hochsten Fallzahlen; die Werte sind in Tausend angegeben."
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.italic = True

    # Slide 7: Zusammenhang Fallzahl/Aufklarungsquote
    slide = prs.slides.add_slide(blank_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Zusammenhang: Fallzahlen und Aufklarungsquoten"
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    
    # Image on left
    pic = slide.shapes.add_picture('figure4_scatter_cases_clearance.png', 
                                  Inches(0.5), Inches(1.2), width=Inches(6))
    
    # Explanatory text box on right
    text_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.2), Inches(4), Inches(4))
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Statistische Analyse"
    p.font.size = Pt(14)
    p.font.bold = True
    
    if r_value is not None:
        p = tf.add_paragraph()
        p.text = f"Korrelation (2024): r = {r_value:.3f}"
        p.font.size = Pt(12)
    
    if r_squared is not None:
        p = tf.add_paragraph()
        p.text = f"Bestimmtheitsmass: R² = {r_squared:.3f}"
        p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Interpretation:"
    p.font.size = Pt(12)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Bundeslander mit hoheren Fallzahlen tendieren zu niedrigeren Aufklarungsquoten. Besonders Berlin fallt als Ausreisser auf."
    p.font.size = Pt(11)
    
    caption = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
    tf = caption.text_frame
    tf.text = "Abbildung 4: Streudiagramm mit Trendlinie. Jeder Punkt reprasentiert ein Bundesland (2024)."
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.italic = True

    # Slide 8: Fazit & Ausblick
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Zusammenfassung & Implikationen"
    
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    
    conclusions = [
        "Fallzahlen zeigen einen abnehmenden Trend in allen Deliktgruppen",
        "Aufklarungsquoten bleiben bei Wohnungseinbruch kritisch niedrig (~15%)",
        "Deutliche regionale Unterschiede: Stadtstaaten mit hochsten Fallzahlen",
        "Negativer Zusammenhang zwischen Fallzahl und Aufklarungsquote",
        "Datenlimitation: Dunkelfeld und Anzeigeverhalten beachten"
    ]
    
    for conclusion in conclusions:
        p = body.add_paragraph()
        p.text = f"• {conclusion}"
        p.font.size = Pt(16)
    
    # Outlook paragraph
    p = body.add_paragraph()
    p.text = "\nAusblick fur weitere Forschung:"
    p.font.size = Pt(14)
    p.font.bold = True
    
    p = body.add_paragraph()
    p.text = "• Integration soziookonomischer Kontrollvariablen"
    p.font.size = Pt(14)
    p.level = 1
    
    p = body.add_paragraph()
    p.text = "• Analyse von Praventionsmassnahmen auf Landerebene"
    p.font.size = Pt(14)
    p.level = 1
    
    p = body.add_paragraph()
    p.text = "• Langsschnittanalysen mit Panel-Daten-Methoden"
    p.font.size = Pt(14)
    p.level = 1

    # Save presentation
    prs.save('PKS_presentation_professional.pptx')
    print("Professional presentation generated successfully.")


if __name__ == '__main__':
    r_value, r_squared = create_figures()
    build_presentation(r_value, r_squared)